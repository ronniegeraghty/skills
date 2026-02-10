"""
Content extraction module for the Executive Review skill.

Handles extraction of text, transcripts, and images from various file types:
- Video/Audio: Whisper transcription, optional frame extraction
- Documents: PDF, DOCX, MD, TXT text and image extraction
- Presentations: PPTX text, notes, and slide images
"""

import io
import os
from datetime import datetime
from pathlib import Path
from typing import Optional

from .constants import (
    DEFAULT_FRAME_INTERVAL_SECONDS,
    MAX_FRAMES_TO_EXTRACT,
    WHISPER_DEFAULT_MODEL,
)
from .types import (
    ContentType,
    ExtractionMetadata,
    ExtractionResult,
    SlideContent,
    TimestampedSegment,
)
from .utils import (
    clean_text,
    ensure_directory,
    format_file_size,
    get_audio_duration,
    get_file_extension,
    get_video_duration,
    print_status,
)


# ============================================================================
# Whisper Transcription
# ============================================================================


def extract_transcript(
    file_path: Path,
    model: str = WHISPER_DEFAULT_MODEL,
    output_dir: Optional[Path] = None,
) -> ExtractionResult:
    """
    Extract transcript from video or audio file using Whisper.

    Args:
        file_path: Path to video or audio file
        model: Whisper model to use (tiny, base, small, medium, large)
        output_dir: Directory to save any output files

    Returns:
        ExtractionResult with transcript text and segments
    """
    import whisper

    print_status(f"Loading Whisper model: {model}", "progress")
    whisper_model = whisper.load_model(model)

    print_status(f"Transcribing: {file_path.name}", "progress")
    result = whisper_model.transcribe(str(file_path))

    # Extract segments with timestamps
    segments = []
    for segment in result.get("segments", []):
        segments.append(
            TimestampedSegment(
                start=segment["start"],
                end=segment["end"],
                text=segment["text"].strip(),
            )
        )

    # Get duration
    duration = get_video_duration(file_path) or get_audio_duration(file_path)

    # Determine content type
    from .detect import detect_content_type

    content_type = detect_content_type(file_path) or ContentType.AUDIO

    # Build metadata
    metadata = ExtractionMetadata(
        file_path=file_path,
        file_name=file_path.name,
        file_size_bytes=file_path.stat().st_size,
        content_type=content_type,
        duration_seconds=duration,
        language=result.get("language"),
    )

    print_status(f"Transcription complete: {len(segments)} segments", "success")

    return ExtractionResult(
        metadata=metadata,
        text=result["text"].strip(),
        segments=segments,
    )


# ============================================================================
# Video Frame Extraction
# ============================================================================


def extract_frames(
    file_path: Path,
    output_dir: Path,
    interval_seconds: int = DEFAULT_FRAME_INTERVAL_SECONDS,
    max_frames: int = MAX_FRAMES_TO_EXTRACT,
) -> list[Path]:
    """
    Extract key frames from video at specified interval.

    Args:
        file_path: Path to video file
        output_dir: Directory to save extracted frames
        interval_seconds: Seconds between frame captures
        max_frames: Maximum number of frames to extract

    Returns:
        List of paths to extracted frame images
    """
    import cv2

    print_status(f"Extracting frames from: {file_path.name}", "progress")

    # Open video
    video = cv2.VideoCapture(str(file_path))
    if not video.isOpened():
        raise ValueError(f"Could not open video file: {file_path}")

    fps = video.get(cv2.CAP_PROP_FPS)
    total_frames = int(video.get(cv2.CAP_PROP_FRAME_COUNT))
    duration = total_frames / fps if fps > 0 else 0

    # Calculate frame interval
    frame_interval = int(fps * interval_seconds)
    if frame_interval < 1:
        frame_interval = 1

    # Ensure output directory exists
    frames_dir = ensure_directory(output_dir / "frames")

    extracted_paths = []
    frame_count = 0
    current_frame = 0

    while True:
        video.set(cv2.CAP_PROP_POS_FRAMES, current_frame)
        success, frame = video.read()

        if not success or frame_count >= max_frames:
            break

        # Calculate timestamp
        timestamp = current_frame / fps if fps > 0 else 0
        timestamp_str = f"{int(timestamp // 60):02d}m{int(timestamp % 60):02d}s"

        # Save frame
        frame_path = frames_dir / f"frame_{frame_count:04d}_{timestamp_str}.jpg"
        cv2.imwrite(str(frame_path), frame)
        extracted_paths.append(frame_path)

        frame_count += 1
        current_frame += frame_interval

    video.release()

    print_status(f"Extracted {len(extracted_paths)} frames", "success")
    return extracted_paths


# ============================================================================
# Document Extraction
# ============================================================================


def extract_pdf(file_path: Path, output_dir: Optional[Path] = None) -> ExtractionResult:
    """
    Extract text and images from PDF file.

    Args:
        file_path: Path to PDF file
        output_dir: Optional directory to save extracted images

    Returns:
        ExtractionResult with text and image paths
    """
    import fitz  # PyMuPDF

    print_status(f"Extracting PDF: {file_path.name}", "progress")

    doc = fitz.open(str(file_path))
    text_parts = []
    image_paths = []

    for page_num, page in enumerate(doc):
        # Extract text
        text_parts.append(f"--- Page {page_num + 1} ---")
        text_parts.append(page.get_text())

        # Extract images if output directory provided
        if output_dir:
            images_dir = ensure_directory(output_dir / "images")
            image_list = page.get_images()

            for img_idx, img in enumerate(image_list):
                try:
                    xref = img[0]
                    pix = fitz.Pixmap(doc, xref)

                    if pix.n - pix.alpha > 3:  # CMYK
                        pix = fitz.Pixmap(fitz.csRGB, pix)

                    img_path = images_dir / f"page{page_num + 1}_img{img_idx + 1}.png"
                    pix.save(str(img_path))
                    image_paths.append(img_path)
                except Exception:
                    pass  # Skip problematic images

    doc.close()

    metadata = ExtractionMetadata(
        file_path=file_path,
        file_name=file_path.name,
        file_size_bytes=file_path.stat().st_size,
        content_type=ContentType.DOCUMENT,
        page_count=len(list(fitz.open(str(file_path)))),
    )

    print_status(f"Extracted {metadata.page_count} pages, {len(image_paths)} images", "success")

    return ExtractionResult(
        metadata=metadata,
        text=clean_text("\n".join(text_parts)),
        image_paths=image_paths,
    )


def extract_docx(file_path: Path) -> ExtractionResult:
    """
    Extract text from Word document.

    Args:
        file_path: Path to DOCX file

    Returns:
        ExtractionResult with text content
    """
    from docx import Document

    print_status(f"Extracting DOCX: {file_path.name}", "progress")

    doc = Document(str(file_path))
    paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]

    metadata = ExtractionMetadata(
        file_path=file_path,
        file_name=file_path.name,
        file_size_bytes=file_path.stat().st_size,
        content_type=ContentType.DOCUMENT,
    )

    print_status(f"Extracted {len(paragraphs)} paragraphs", "success")

    return ExtractionResult(
        metadata=metadata,
        text=clean_text("\n\n".join(paragraphs)),
    )


def extract_markdown(file_path: Path) -> ExtractionResult:
    """
    Extract text from Markdown file.

    Args:
        file_path: Path to MD file

    Returns:
        ExtractionResult with text content
    """
    print_status(f"Reading Markdown: {file_path.name}", "progress")

    with open(file_path, "r", encoding="utf-8") as f:
        text = f.read()

    metadata = ExtractionMetadata(
        file_path=file_path,
        file_name=file_path.name,
        file_size_bytes=file_path.stat().st_size,
        content_type=ContentType.DOCUMENT,
    )

    return ExtractionResult(
        metadata=metadata,
        text=text,
    )


def extract_text_file(file_path: Path) -> ExtractionResult:
    """
    Extract text from plain text file.

    Args:
        file_path: Path to TXT file

    Returns:
        ExtractionResult with text content
    """
    print_status(f"Reading text file: {file_path.name}", "progress")

    with open(file_path, "r", encoding="utf-8") as f:
        text = f.read()

    metadata = ExtractionMetadata(
        file_path=file_path,
        file_name=file_path.name,
        file_size_bytes=file_path.stat().st_size,
        content_type=ContentType.DOCUMENT,
    )

    return ExtractionResult(
        metadata=metadata,
        text=text,
    )


def extract_document(
    file_path: Path, output_dir: Optional[Path] = None
) -> ExtractionResult:
    """
    Extract content from a document file (PDF, DOCX, MD, TXT).

    Args:
        file_path: Path to document file
        output_dir: Optional directory for extracted images

    Returns:
        ExtractionResult with text and optional images
    """
    ext = get_file_extension(file_path)

    if ext == ".pdf":
        return extract_pdf(file_path, output_dir)
    elif ext == ".docx":
        return extract_docx(file_path)
    elif ext == ".md":
        return extract_markdown(file_path)
    elif ext == ".txt":
        return extract_text_file(file_path)
    else:
        raise ValueError(f"Unsupported document format: {ext}")


# ============================================================================
# Presentation Extraction
# ============================================================================


def extract_pptx(
    file_path: Path, output_dir: Optional[Path] = None
) -> ExtractionResult:
    """
    Extract text, notes, and slide images from PowerPoint file.

    Args:
        file_path: Path to PPTX file
        output_dir: Optional directory to save slide images

    Returns:
        ExtractionResult with slides and optional images
    """
    from pptx import Presentation
    from pptx.util import Inches

    print_status(f"Extracting PPTX: {file_path.name}", "progress")

    prs = Presentation(str(file_path))
    slides_content = []
    text_parts = []
    image_paths = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        # Extract title
        title = None
        if slide.shapes.title:
            title = slide.shapes.title.text

        # Extract all text from shapes
        slide_text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text_parts.append(shape.text)

        slide_text = "\n".join(slide_text_parts)

        # Extract notes
        notes = None
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            if notes_text_frame and notes_text_frame.text.strip():
                notes = notes_text_frame.text.strip()

        # Create slide content
        slide_content = SlideContent(
            slide_number=slide_num,
            title=title,
            text=clean_text(slide_text),
            notes=notes,
        )
        slides_content.append(slide_content)

        # Build text representation
        text_parts.append(f"--- Slide {slide_num}: {title or 'Untitled'} ---")
        text_parts.append(slide_text)
        if notes:
            text_parts.append(f"[Speaker Notes: {notes}]")
        text_parts.append("")

    # Note: python-pptx doesn't directly export slide images
    # Would need to use pdf2image or similar for that
    # For now, we just extract text content

    metadata = ExtractionMetadata(
        file_path=file_path,
        file_name=file_path.name,
        file_size_bytes=file_path.stat().st_size,
        content_type=ContentType.PRESENTATION,
        slide_count=len(slides_content),
    )

    print_status(f"Extracted {len(slides_content)} slides", "success")

    return ExtractionResult(
        metadata=metadata,
        text="\n".join(text_parts),
        slides=slides_content,
        image_paths=image_paths,
    )


def extract_presentation(
    file_path: Path, output_dir: Optional[Path] = None
) -> ExtractionResult:
    """
    Extract content from a presentation file.

    Args:
        file_path: Path to presentation file
        output_dir: Optional directory for extracted images

    Returns:
        ExtractionResult with slides and optional images
    """
    ext = get_file_extension(file_path)

    if ext == ".pptx":
        return extract_pptx(file_path, output_dir)
    elif ext == ".pdf":
        # Treat PDF as presentation with images
        return extract_pdf(file_path, output_dir)
    else:
        raise ValueError(f"Unsupported presentation format: {ext}")


# ============================================================================
# Main Extraction Dispatcher
# ============================================================================


def extract_content(
    file_path: Path,
    content_type: ContentType,
    output_dir: Optional[Path] = None,
    enable_frames: bool = False,
    whisper_model: str = WHISPER_DEFAULT_MODEL,
    frame_interval: int = DEFAULT_FRAME_INTERVAL_SECONDS,
) -> ExtractionResult:
    """
    Main extraction dispatcher. Routes to appropriate extractor based on content type.

    Args:
        file_path: Path to the file to extract
        content_type: Type of content
        output_dir: Optional directory for output files
        enable_frames: Whether to extract video frames (for video only)
        whisper_model: Whisper model to use for transcription
        frame_interval: Seconds between frame captures

    Returns:
        ExtractionResult with extracted content
    """
    if content_type == ContentType.VIDEO:
        # Transcribe video
        result = extract_transcript(file_path, model=whisper_model, output_dir=output_dir)

        # Optionally extract frames
        if enable_frames and output_dir:
            frame_paths = extract_frames(
                file_path, output_dir, interval_seconds=frame_interval
            )
            result.image_paths = frame_paths

        return result

    elif content_type == ContentType.AUDIO:
        return extract_transcript(file_path, model=whisper_model, output_dir=output_dir)

    elif content_type == ContentType.DOCUMENT:
        return extract_document(file_path, output_dir)

    elif content_type == ContentType.PRESENTATION:
        return extract_presentation(file_path, output_dir)

    else:
        raise ValueError(f"Unknown content type: {content_type}")
